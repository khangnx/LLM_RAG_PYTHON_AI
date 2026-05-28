"""
main.py
-------
FastAPI Backend cho hệ thống AI RAG Demo.
Hỗ trợ đọc nhiều định dạng file: PDF, Excel, Word, TXT, LOG, PowerPoint.
Lưu lịch sử chat vào MySQL thông qua SQLAlchemy.
"""

import os
import time
import logging
from datetime import datetime
from fastapi import APIRouter, HTTPException
import chromadb

from fastapi import FastAPI, HTTPException, Depends
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import pandas as pd

# LangChain components
from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_community.vectorstores import Chroma
from langchain_community.llms import Ollama
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser

# Thêm module cho trí nhớ hội thoại (Conversational Memory)
from langchain.chains import create_history_aware_retriever, create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_community.chat_message_histories import ChatMessageHistory
from langchain_core.chat_history import BaseChatMessageHistory
from langchain_core.runnables.history import RunnableWithMessageHistory

# Embedding model (HuggingFace local, không cần API key)
from langchain_community.embeddings import HuggingFaceEmbeddings

# Thư viện đọc file Word và PowerPoint
import docx
from pptx import Presentation

# SQLAlchemy session
from sqlalchemy.orm import Session

# Module Database nội bộ
from database import init_db, get_db, RagHistory, ChatSession, ChatMessage

# =====================================================
# KHỞI TẠO LOGGING
# =====================================================
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# =====================================================
# KHỞI TẠO EMBEDDING MODEL
# =====================================================
embeddings = HuggingFaceEmbeddings(model_name="paraphrase-multilingual-MiniLM-L12-v2")

# =====================================================
# KHỞI TẠO FASTAPI APP
# =====================================================
app = FastAPI(title="AI RAG Demo API")

# Cấu hình CORS để VueJS từ máy Host gọi vào container được
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# =====================================================
# ĐƯỜNG DẪN THƯ MỤC DỮ LIỆU
# =====================================================
DATA_SOURCE_DIR = "/data_source"
VECTOR_DB_DIR = "/vector_db"

# =====================================================
# KHỞI TẠO LLM OLLAMA
# =====================================================
OLLAMA_HOST = os.getenv("OLLAMA_HOST", "http://ollama:11434")
llm = Ollama(base_url=OLLAMA_HOST, model="llama3")


# =====================================================
# LIFECYCLE EVENT - KHỞI TẠO DATABASE KHI APP START
# =====================================================
@app.on_event("startup")
def on_startup():
    """
    Hàm chạy khi FastAPI khởi động.
    Thử kết nối MySQL và tạo bảng (nếu chưa có).
    Retry tối đa 10 lần, mỗi lần chờ 3 giây (MySQL cần thời gian khởi động).
    """
    max_retries = 10
    for attempt in range(max_retries):
        try:
            init_db()
            logger.info("✅ Kết nối MySQL và khởi tạo bảng thành công!")
            return
        except Exception as e:
            logger.warning(f"⏳ MySQL chưa sẵn sàng (lần {attempt + 1}/{max_retries}): {e}")
            time.sleep(3)
    logger.error("❌ Không thể kết nối MySQL sau nhiều lần thử. Chức năng lưu lịch sử sẽ bị tắt.")


# =====================================================
# PYDANTIC MODELS
# =====================================================
class ChatRequest(BaseModel):
    question: str
    session_id: str  # BẮT BUỘC — Frontend phải truyền session_id từ DB

class CreateSessionRequest(BaseModel):
    title: str = "Cuộc trò chuyện mới"

# =====================================================
# HÀM TẢI LỊCH SỬ CHAT TỪ DB (DB-BACKED MEMORY)
# Thay thế in-memory store cũ — lịch sử được lưu vĩnh viễn trong MySQL
# =====================================================
def build_chat_history_from_db(session_id: str, db: Session) -> BaseChatMessageHistory:
    """
    Đọc toàn bộ tin nhắn cũ của session từ MySQL,
    nạp vào ChatMessageHistory để LangChain dùng làm ngữ cảnh.
    """
    history = ChatMessageHistory()
    messages = (
        db.query(ChatMessage)
        .filter(ChatMessage.session_id == session_id)
        .order_by(ChatMessage.created_at.asc())
        .all()
    )
    for msg in messages:
        if msg.role == "user":
            history.add_user_message(msg.content)
        else:
            history.add_ai_message(msg.content)
    return history


# =====================================================
# CÁC HÀM XỬ LÝ TRÍCH XUẤT & CHUNKING TỪNG LOẠI FILE
# =====================================================

# Bộ cắt nhỏ văn bản dùng chung cho tất cả loại file
text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)


def process_pdf(file_path: str) -> list:
    """
    Đọc file PDF, cắt nhỏ từng trang thành các chunk văn bản.
    Gán metadata: source_type=pdf, file_name, số trang.
    """
    try:
        loader = PyPDFLoader(file_path)
        pages = loader.load()
        chunks = text_splitter.split_documents(pages)
        for chunk in chunks:
            chunk.metadata["source_type"] = "pdf"
            chunk.metadata["file_name"] = os.path.basename(file_path)
        logger.info(f"  ✅ PDF '{os.path.basename(file_path)}': {len(chunks)} chunks")
        return chunks
    except Exception as e:
        logger.error(f"  ❌ Lỗi đọc PDF '{file_path}': {e}")
        return []


def process_excel(file_path: str) -> list:
    """
    Đọc file Excel (.xlsx, .xls).
    Mỗi dòng dữ liệu trong mỗi sheet được chuyển thành một Document riêng.
    """
    chunks = []
    file_name = os.path.basename(file_path)
    try:
        excel_file = pd.ExcelFile(file_path)
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name).fillna("")
            for index, row in df.iterrows():
                row_parts = [f"{col}: {row[col]}" for col in df.columns]
                row_text = " | ".join(row_parts)
                doc = Document(
                    page_content=row_text,
                    metadata={
                        "source_type": "excel",
                        "file_name": file_name,
                        "sheet_name": sheet_name,
                        "row_index": index + 1
                    }
                )
                chunks.append(doc)
        logger.info(f"  ✅ Excel '{file_name}': {len(chunks)} chunks")
    except Exception as e:
        logger.error(f"  ❌ Lỗi đọc Excel '{file_path}': {e}")
    return chunks


def process_word(file_path: str) -> list:
    """
    Đọc file Word (.docx) bằng thư viện python-docx.
    Gom toàn bộ các đoạn văn (paragraph) thành một chuỗi lớn, sau đó cắt nhỏ.
    """
    file_name = os.path.basename(file_path)
    try:
        doc = docx.Document(file_path)
        # Gom tất cả đoạn văn, loại bỏ dòng trống
        full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        if not full_text:
            logger.warning(f"  ⚠️ Word '{file_name}' không có nội dung văn bản.")
            return []
        # Tạo Document tổng, sau đó cắt nhỏ
        base_doc = Document(
            page_content=full_text,
            metadata={"source_type": "word", "file_name": file_name}
        )
        chunks = text_splitter.split_documents([base_doc])
        logger.info(f"  ✅ Word '{file_name}': {len(chunks)} chunks")
        return chunks
    except Exception as e:
        logger.error(f"  ❌ Lỗi đọc Word '{file_path}': {e}")
        return []


def process_text(file_path: str, source_type: str = "txt") -> list:
    """
    Đọc file văn bản thuần (.txt) hoặc file log (.log).
    Dùng TextLoader với encoding UTF-8.
    """
    file_name = os.path.basename(file_path)
    try:
        loader = TextLoader(file_path, encoding="utf-8")
        docs = loader.load()
        chunks = text_splitter.split_documents(docs)
        for chunk in chunks:
            chunk.metadata["source_type"] = source_type
            chunk.metadata["file_name"] = file_name
        logger.info(f"  ✅ Text/Log '{file_name}': {len(chunks)} chunks")
        return chunks
    except Exception as e:
        logger.error(f"  ❌ Lỗi đọc Text/Log '{file_path}': {e}")
        return []


def process_pptx(file_path: str) -> list:
    """
    Đọc file PowerPoint (.pptx) bằng thư viện python-pptx.
    Duyệt qua từng Slide, trích xuất text từ toàn bộ Shape/TextBox.
    Gom text của mỗi slide thành một đoạn, sau đó cắt nhỏ toàn bộ.
    """
    file_name = os.path.basename(file_path)
    try:
        prs = Presentation(file_path)
        all_slides_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                # Chỉ xử lý các shape có chứa khung văn bản
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                # Gộp các text trong một slide lại, thêm nhãn số slide
                slide_content = f"[Slide {slide_num}]\n" + "\n".join(slide_texts)
                all_slides_text.append(slide_content)

        full_text = "\n\n".join(all_slides_text)
        if not full_text:
            logger.warning(f"  ⚠️ PowerPoint '{file_name}' không có nội dung.")
            return []

        base_doc = Document(
            page_content=full_text,
            metadata={"source_type": "pptx", "file_name": file_name}
        )
        chunks = text_splitter.split_documents([base_doc])
        logger.info(f"  ✅ PowerPoint '{file_name}': {len(chunks)} chunks từ {len(prs.slides)} slides")
        return chunks
    except Exception as e:
        logger.error(f"  ❌ Lỗi đọc PowerPoint '{file_path}': {e}")
        return []
def process_log_smart(file_path: str) -> list:
    """
    Xử lý file .log thông minh:
    - Nếu nội dung là JSON array ([ {...}, {...} ]) → parse từng object thành Document riêng,
      mỗi trường key-value được viết rõ ràng dạng "key: value" để AI dễ tìm kiếm.
    - Nếu là log văn bản thông thường → fallback về TextLoader như cũ.
    """
    import json
    file_name = os.path.basename(file_path)
    chunks = []

    try:
        with open(file_path, "r", encoding="utf-8") as f:
            raw = f.read().strip()

        # Thử parse như JSON array
        try:
            data = json.loads(raw)
            if isinstance(data, list) and len(data) > 0 and isinstance(data[0], dict):
                logger.info(f"  📋 Log '{file_name}' phát hiện JSON array ({len(data)} records) — xử lý như bảng dữ liệu.")
                for index, record in enumerate(data):
                    # Chuyển từng object JSON thành dạng "key: value | key: value" để AI đọc hiểu
                    row_parts = [f"{k}: {v}" for k, v in record.items() if v is not None]
                    row_text = " | ".join(row_parts)
                    doc = Document(
                        page_content=row_text,
                        metadata={
                            "source_type": "log_json",
                            "file_name": file_name,
                            "row_index": index + 1
                        }
                    )
                    chunks.append(doc)
                logger.info(f"  ✅ Log JSON '{file_name}': {len(chunks)} chunks (1 chunk/record)")
                return chunks
        except json.JSONDecodeError:
            pass  # Không phải JSON → xử lý như text thuần bên dưới

        # Fallback: đọc như text thuần
        logger.info(f"  📄 Log '{file_name}' là văn bản thuần — dùng TextLoader.")
        return process_text(file_path, source_type="log")

    except Exception as e:
        logger.error(f"  ❌ Lỗi đọc Log '{file_path}': {e}")
        return []




@app.get("/api/files")
def get_files():
    """
    Lấy danh sách tất cả file tài liệu đang có trong thư mục data_source.
    Hỗ trợ: .pdf, .xlsx, .xls, .docx, .txt, .log, .pptx
    """
    supported_extensions = ('.pdf', '.xlsx', '.xls', '.docx', '.txt', '.log', '.pptx')
    if not os.path.exists(DATA_SOURCE_DIR):
        return {"files": []}
    files = [
        f for f in os.listdir(DATA_SOURCE_DIR)
        if f.lower().endswith(supported_extensions)
    ]
    return {"files": sorted(files)}


@app.post("/api/ingest")
def ingest_data():
    """
    Quét toàn bộ thư mục data_source, trích xuất văn bản từ mọi loại file,
    cắt nhỏ thành chunk và lưu vector vào ChromaDB.
    Bọc try/except từng file để tránh 1 file lỗi làm chết cả API.
    """
    if not os.path.exists(DATA_SOURCE_DIR):
        raise HTTPException(status_code=404, detail="Thư mục data_source không tồn tại.")

    all_chunks = []
    supported_extensions = ('.pdf', '.xlsx', '.xls', '.docx', '.txt', '.log', '.pptx')

    logger.info(f"🔄 Bắt đầu Ingest dữ liệu từ '{DATA_SOURCE_DIR}'...")

    for file_name in os.listdir(DATA_SOURCE_DIR):
        file_path = os.path.join(DATA_SOURCE_DIR, file_name)
        lower_name = file_name.lower()

        if lower_name.endswith('.pdf'):
            all_chunks.extend(process_pdf(file_path))
        elif lower_name.endswith(('.xlsx', '.xls')):
            all_chunks.extend(process_excel(file_path))
        elif lower_name.endswith('.docx'):
            all_chunks.extend(process_word(file_path))
        elif lower_name.endswith('.txt'):
            all_chunks.extend(process_text(file_path, source_type="txt"))
        elif lower_name.endswith('.log'):
            all_chunks.extend(process_log_smart(file_path))
        elif lower_name.endswith('.pptx'):
            all_chunks.extend(process_pptx(file_path))

    if not all_chunks:
        return {
            "status": "warning",
            "message": "Không tìm thấy file hợp lệ hoặc không có dữ liệu nào được trích xuất.",
            "chunks": 0
        }

    # Làm sạch và lọc bỏ các chunk không hợp lệ
    # - Re-encode UTF-8 để loại ký tự không hợp lệ từ PDF/Word (toán học, LaTeX, byte lạ)
    # - Xóa null byte \x00 và control characters
    # - Bỏ qua chunk quá ngắn (< 10 ký tự) hoặc chỉ chứa dấu phân tách
    cleaned_chunks = []
    for doc in all_chunks:
        if not doc or not hasattr(doc, "page_content"):
            continue
        raw = doc.page_content if doc.page_content is not None else ""
        # Re-encode UTF-8 để strip bytes không hợp lệ
        raw = raw.encode("utf-8", errors="ignore").decode("utf-8", errors="ignore")
        # Xóa null byte và ký tự control ASCII (0x00-0x08, 0x0B, 0x0C, 0x0E-0x1F)
        raw = "".join(c for c in raw if c == "\n" or c == "\t" or ord(c) >= 32)
        raw = raw.strip()
        # Bỏ qua nếu quá ngắn hoặc chỉ toàn dấu phân tách
        if len(raw) < 10:
            continue
        if all(c in ' |:\n\t\r-_.,' for c in raw):
            continue
        doc.page_content = raw
        cleaned_chunks.append(doc)

    all_chunks = cleaned_chunks

    if not all_chunks:
        return {
            "status": "warning",
            "message": "Sau khi làm sạch, không còn dữ liệu văn bản hợp lệ nào để lưu.",
            "chunks": 0
        }

    # Debug log để kiểm tra kiểu dữ liệu của page_content trong từng chunk
    logger.info(f"🔍 Kiểm tra dữ liệu: Tổng số {len(all_chunks)} chunks.")
    for idx, doc in enumerate(all_chunks):
        if not isinstance(doc.page_content, str):
            logger.error(f"❌ CHUNK #{idx} KHÔNG PHẢI STRING: {type(doc.page_content)} (metadata: {doc.metadata})")
        elif not doc.page_content.strip():
            logger.error(f"❌ CHUNK #{idx} BỊ RỖNG: (metadata: {doc.metadata})")

    # Xóa nội dung VectorDB cũ an toàn bằng API của Chroma thay vì xóa file trực tiếp
    try:
        old_store = Chroma(persist_directory=VECTOR_DB_DIR, embedding_function=embeddings)
        old_store.delete_collection()
        logger.info("🗑️ Đã dọn dẹp collection VectorDB cũ thành công.")
    except Exception as e:
        logger.warning(f"⚠️ Không thể xóa collection cũ (có thể chưa tồn tại): {e}")


    # Lưu toàn bộ vector vào ChromaDB
    try:
        vector_store = Chroma.from_documents(
            documents=all_chunks,
            embedding=embeddings,
            persist_directory=VECTOR_DB_DIR
        )
        vector_store.persist()
        logger.info(f"✅ Đã lưu {len(all_chunks)} chunks vào VectorDB.")
    except Exception as e:
        logger.exception("❌ Lỗi khi lưu vào VectorDB:")
        raise HTTPException(status_code=500, detail=f"Lỗi khi lưu vào VectorDB: {str(e)}")

    return {
        "status": "success",
        "message": f"Đồng bộ thành công! Đã xử lý {len(all_chunks)} chunks.",
        "chunks": len(all_chunks)
    }


# =====================================================
# API: QUẢN LÝ PHIÊN HỘI THOẠI (SESSIONS)
# =====================================================

@app.post("/api/sessions", status_code=201)
def create_session(payload: CreateSessionRequest, db: Session = Depends(get_db)):
    """
    Tạo một cuộc trò chuyện mới. Tự động sinh UUID làm khóa chính.
    Frontend gọi API này khi user nhấn '+ Cuộc trò chuyện mới'.
    """
    new_session = ChatSession(title=payload.title)
    db.add(new_session)
    db.commit()
    db.refresh(new_session)
    logger.info(f"✅ Tạo session mới: {new_session.id}")
    return {
        "id": new_session.id,
        "title": new_session.title,
        "created_at": new_session.created_at.strftime("%Y-%m-%d %H:%M:%S")
    }


@app.get("/api/sessions")
def list_sessions(db: Session = Depends(get_db)):
    """
    Lấy tất cả phiên hội thoại, sắp xếp mới nhất lên đầu.
    Dùng để render danh sách sidebar bên trái.
    """
    sessions = db.query(ChatSession).order_by(ChatSession.created_at.desc()).all()
    return [
        {
            "id": s.id,
            "title": s.title,
            "created_at": s.created_at.strftime("%Y-%m-%d %H:%M:%S") if s.created_at else ""
        }
        for s in sessions
    ]


@app.get("/api/sessions/{session_id}/messages")
def get_session_messages(session_id: str, db: Session = Depends(get_db)):
    """
    Tải toàn bộ lịch sử tin nhắn của một phiên cụ thể.
    Frontend gọi API này khi user click vào một item trên Sidebar.
    """
    # Kiểm tra session tồn tại
    session = db.query(ChatSession).filter(ChatSession.id == session_id).first()
    if not session:
        raise HTTPException(status_code=404, detail="Phiên hội thoại không tồn tại.")

    messages = (
        db.query(ChatMessage)
        .filter(ChatMessage.session_id == session_id)
        .order_by(ChatMessage.created_at.asc())
        .all()
    )
    return {
        "session_id": session_id,
        "title": session.title,
        "messages": [
            {
                "id": m.id,
                "role": m.role,
                "content": m.content,
                "created_at": m.created_at.strftime("%Y-%m-%d %H:%M:%S") if m.created_at else ""
            }
            for m in messages
        ]
    }


@app.post("/api/chat")
def chat_with_rag(request: ChatRequest, db: Session = Depends(get_db)):
    """
    Tiếp nhận câu hỏi từ user, truy vấn VectorDB để lấy ngữ cảnh liên quan,
    ráp Prompt hoàn chỉnh gửi sang Ollama, nhận câu trả lời và lưu lịch sử vào MySQL.
    """
    # Kiểm tra VectorDB đã được Ingest chưa
    if not os.path.exists(VECTOR_DB_DIR) or not os.listdir(VECTOR_DB_DIR):
        raise HTTPException(
            status_code=400,
            detail="Vui lòng nhấn nút 'Đồng bộ hệ thống' (Ingest) trước khi chat."
        )

    # Kiểm tra session có tồn tại không, nếu không tạo mới
    session_obj = db.query(ChatSession).filter(ChatSession.id == request.session_id).first()
    if not session_obj:
        raise HTTPException(status_code=404, detail="Phiên hội thoại không tồn tại. Hãy tạo mới.")

    try:
        # Kết nối VectorDB và tìm 5 đoạn tài liệu liên quan nhất
        vector_store = Chroma(persist_directory=VECTOR_DB_DIR, embedding_function=embeddings)
        retriever = vector_store.as_retriever(search_kwargs={"k": 10})
        relevant_docs = retriever.invoke(request.question)

        # Bổ sung: tìm kiếm thêm bằng keyword exact-match (chỉ áp dụng với mã sản phẩm/từ khóa kỹ thuật)
        # Regex: Từ có chứa ít nhất 1 chữ số (ví dụ SP001, 100m2) HOẶC từ in hoa hoàn toàn >= 3 ký tự (VND, VAT)
        import re
        from langchain_core.documents import Document
        meaningful_tokens = re.findall(r'[a-zA-Z_]*\d+[a-zA-Z0-9_]*|[A-Z_]{3,}', request.question)
        if meaningful_tokens:
            keyword_query = meaningful_tokens[0] # Lấy từ khóa quan trọng nhất để tìm exact match
            logger.info(f"🔑 Keyword exact match với: '{keyword_query}'")
            try:
                # Sử dụng API nội bộ của Chroma để tìm chính xác văn bản chứa keyword
                exact_results = vector_store._collection.get(
                    where_document={"$contains": keyword_query},
                    include=["documents", "metadatas"]
                )
                
                # Dedup dựa trên nội dung thực tế (hash) thay vì id() object
                existing_contents = {doc.page_content for doc in relevant_docs}
                
                if exact_results and exact_results["documents"]:
                    added_exact = 0
                    for idx, doc_text in enumerate(exact_results["documents"]):
                        if added_exact >= 10:  # Giới hạn tối đa 10 chunk exact-match
                            break
                        if doc_text not in existing_contents:
                            # Chuyển đổi thành Document object để tương thích với luồng RAG
                            meta = exact_results["metadatas"][idx] if exact_results["metadatas"] else {}
                            relevant_docs.append(Document(page_content=doc_text, metadata=meta))
                            existing_contents.add(doc_text)
                            added_exact += 1
                            logger.info(f"✅ Đã thêm document exact-match cho keyword '{keyword_query}'")
            except Exception as e:
                logger.warning(f"⚠️ Lỗi khi tìm kiếm exact match: {e}")

        # ---------------------------------------------------------
        # TẠO HISTORY-AWARE RETRIEVER (Hiểu ngữ cảnh từ lịch sử chat)
        # ---------------------------------------------------------
        contextualize_q_system_prompt = (
            "Dựa trên lịch sử trò chuyện và câu hỏi mới nhất của người dùng, "
            "có thể người dùng đang tham chiếu đến ngữ cảnh trong lịch sử trò chuyện. "
            "Hãy viết lại một câu hỏi độc lập sao cho nó có thể tự hiểu được mà không cần lịch sử trò chuyện. "
            "KHÔNG ĐƯỢC trả lời câu hỏi, chỉ viết lại nó nếu cần, nếu không cần thì giữ nguyên."
        )
        contextualize_q_prompt = ChatPromptTemplate.from_messages([
            ("system", contextualize_q_system_prompt),
            MessagesPlaceholder("chat_history"),
            ("human", "{input}"),
        ])
        
        history_aware_retriever = create_history_aware_retriever(
            llm, retriever, contextualize_q_prompt
        )

        # ---------------------------------------------------------
        # TẠO QA CHAIN (Ráp context vào prompt cuối)
        # ---------------------------------------------------------
        qa_system_prompt = """Bạn là một trợ lý AI thông minh chuyên phân tích dữ liệu nội bộ.
YÊU CẦU QUAN TRỌNG:
1. LUÔN LUÔN TRẢ LỜI BẰNG TIẾNG VIỆT (VIETNAMESE).
2. Hãy trả lời câu hỏi dựa TRỰC TIẾP vào phần ngữ cảnh (Context) bên dưới.
3. Nếu ngữ cảnh là một câu chuyện hoặc đoạn mô tả, hãy tổng hợp các thông tin đó để trả lời thay vì đòi hỏi định nghĩa chính xác.
4. Nếu thông tin không có trong ngữ cảnh, hãy trả lời chính xác câu này: 'Tôi không tìm thấy thông tin này trong tài liệu của bạn'.
5. Tuyệt đối không tự bịa câu trả lời ngoài ngữ cảnh.

Ngữ cảnh (Context):
{context}"""
        
        qa_prompt = ChatPromptTemplate.from_messages([
            ("system", qa_system_prompt),
            MessagesPlaceholder("chat_history"),
            ("human", "{input}"),
        ])

        question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)
        rag_chain = create_retrieval_chain(history_aware_retriever, question_answer_chain)

        # Tải lịch sử chat từ DB (thay in-memory store)
        db_history = build_chat_history_from_db(request.session_id, db)

        # Gắn Message History theo session — dùng lambda để giữ db_history đã load
        conversational_rag_chain = RunnableWithMessageHistory(
            rag_chain,
            lambda sid: db_history,
            input_messages_key="input",
            history_messages_key="chat_history",
            output_messages_key="answer",
        )

        # Thực thi chain
        response = conversational_rag_chain.invoke(
            {"input": request.question},
            config={"configurable": {"session_id": request.session_id}}
        )
        
        answer = response["answer"]
        # context được trả về từ history_aware_retriever
        retrieved_docs = response["context"]

        # ---------------------------------------------------------
        # RENDER PROMPT LOG (để hiển thị cho khách hàng xem)
        # ---------------------------------------------------------
        # Lấy context cuối cùng mà chain đã dùng để trả lời
        context_text = "\n\n".join([doc.page_content for doc in retrieved_docs])
        
        # Nếu relevant_docs có exact match, ta add thêm vào retrieved_docs để tính nguồn
        for d in relevant_docs:
            if d not in retrieved_docs:
                retrieved_docs.append(d)
                
        rendered_prompt = qa_system_prompt.replace("{context}", context_text) + "\n\nCâu hỏi: " + request.question

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi khi xử lý câu hỏi: {str(e)}")

    # Lưu tin nhắn vào bảng chat_messages và cập nhật audit log
    try:
        # 1. Lưu tin nhắn user
        db.add(ChatMessage(session_id=request.session_id, role="user", content=request.question))
        # 2. Lưu tin nhắn AI
        db.add(ChatMessage(session_id=request.session_id, role="ai", content=answer))
        # 3. Nếu đây là tin nhắn đầu tiên, cập nhật title của session theo câu hỏi
        msg_count = db.query(ChatMessage).filter(ChatMessage.session_id == request.session_id).count()
        if msg_count == 0 and session_obj.title == "Cuộc trò chuyện mới":
            session_obj.title = request.question[:80]  # Cắt tối đa 80 ký tự
            db.add(session_obj)
        # 4. Lưu audit log vào rag_history (giữ nguyên tính năng cũ)
        db.add(RagHistory(
            search_query=request.question,
            generated_prompt=rendered_prompt,
            ai_response=answer,
            created_at=datetime.utcnow()
        ))
        db.commit()
        logger.info(f"✅ Đã lưu tin nhắn vào session '{request.session_id[:8]}...' và audit log.")
    except Exception as e:
        logger.error(f"⚠️ Không thể lưu lịch sử vào MySQL: {e}")
        db.rollback()

    # Trích xuất thông tin nguồn trích dẫn (Citations) gửi về cho VueJS hiển thị
    sources = []
    for doc in retrieved_docs:
        meta = doc.metadata
        source_info = {"file": meta.get("file_name", "Unknown")}

        if meta.get("source_type") == "pdf":
            source_info["location"] = f"Trang {meta.get('page', 0) + 1}"
        elif meta.get("source_type") == "excel":
            source_info["location"] = f"Sheet: {meta.get('sheet_name')}, Dòng: {meta.get('row_index')}"
        elif meta.get("source_type") == "pptx":
            source_info["location"] = "Slide (xem file để biết thêm chi tiết)"
        elif meta.get("source_type") in ("txt", "log"):
            source_info["location"] = "Văn bản/Log"
        elif meta.get("source_type") == "word":
            source_info["location"] = "Tài liệu Word"
        else:
            source_info["location"] = "Không rõ"

        # Tránh thêm nguồn trùng lặp
        if source_info not in sources:
            sources.append(source_info)

    return {
        "answer": answer,
        "sources": sources
    }


@app.get("/api/history")
def get_history(db: Session = Depends(get_db)):
    """
    Lấy toàn bộ lịch sử chat từ bảng rag_history trong MySQL.
    Sắp xếp theo thời gian mới nhất lên đầu.
    """
    try:
        records = db.query(RagHistory).order_by(RagHistory.created_at.desc()).all()
        return [
            {
                "id": r.id,
                "search_query": r.search_query,
                "generated_prompt": r.generated_prompt,
                "ai_response": r.ai_response,
                "created_at": r.created_at.strftime("%Y-%m-%d %H:%M:%S") if r.created_at else ""
            }
            for r in records
        ]
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi khi lấy lịch sử: {str(e)}")

# Khởi tạo client kết nối trực tiếp tới phân vùng Docker Volume /vector_db
# (Giống hệt cách file diagnose.py đã kết nối thành công)
try:
    chroma_client = chromadb.PersistentClient(path="/vector_db")
except Exception as e:
    chroma_client = None
# Thêm endpoint này vào dưới các endpoint /api/chat hoặc /api/files của bạn
@app.get("/api/collections")
def get_chroma_collections():
    if not chroma_client:
        raise HTTPException(status_code=500, detail="ChromaDB client không sẵn sàng")
    
    try:
        collections = chroma_client.list_collections()
        result = []
        for col in collections:
            collection_obj = chroma_client.get_collection(name=col.name)
            chunk_count = collection_obj.count()
            
            # LẤY THÊM: 5 đoạn văn bản mẫu đầu tiên để xem thử nội dung
            sample_docs = []
            if chunk_count > 0:
                sample = collection_obj.get(limit=50)
                sample_docs = sample.get('documents', [])

            result.append({
                "name": col.name,
                "id": str(col.id),
                "count": chunk_count,
                "metadata": col.metadata,
                "samples": sample_docs  # Trả thêm danh sách mẫu về cho Frontend UI
            })
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))